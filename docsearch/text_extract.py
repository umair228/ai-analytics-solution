"""
Document text extraction for the AI Document Search feature.

Ported and extended from the LW-Chatbot ``upload_processor.py``. Walks a corpus
directory and turns every supported file into sentence-level *chunks* with
provenance metadata, ready for the BM25 / semantic index in ``retrieval.py``.

Supported types (extensions, case-insensitive):
    .docx           Word documents  (python-docx)  -- the real LIMS corpus
    .pptx           PowerPoint       (python-pptx)
    .pdf            PDF              (pypdf)
    .txt .md        Plain text
    .png .jpg .jpeg Images, OCR'd    (easyocr -- optional, off by default)

Every extractor is wrapped so a single unreadable file (e.g. a git-LFS pointer
stub left un-hydrated) is skipped with a warning rather than aborting the build.

Public API
----------
    extract_corpus(corpus_dir) -> (chunk_texts: list[str], chunk_metas: list[dict])

``chunk_metas`` entries are ``{"source": filename, "page": int|str, "sentence": str}``.
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Dict, List, Tuple

SUPPORTED_EXTS = {".docx", ".pptx", ".pdf", ".txt", ".md", ".png", ".jpg", ".jpeg"}
TEXT_EXTS = {".txt", ".md"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg"}

_SENT_SPLIT_RE = re.compile(r"(?<=[.!?])\s+")
_MIN_SENT_CHARS = 5


def _split_sentences(text: str) -> List[str]:
    """Split text into sentence-ish chunks, dropping very short fragments."""
    if not text:
        return []
    return [s.strip() for s in _SENT_SPLIT_RE.split(text) if s and len(s.strip()) >= _MIN_SENT_CHARS]


# ──────────────────────────────────────────────────────────────────────────
# Passage packing
#
# The corpus used to be chunked one *sentence* per chunk (~7 tokens), which is
# far too small for either BM25 or dense retrieval to be meaningful. Instead we
# pack consecutive sentences into ~``DOCSEARCH_CHUNK_TOKENS``-token *passages*
# (hard-capped, with a small sentence overlap so context isn't sheared at the
# boundary). The packer is intentionally dependency-free (a word-count token
# estimate, no tokenizer) so it never pulls a numpy>=2 transitive dependency.
# ``_est_tokens`` and ``pack_passages`` are reused by the knowledge-base
# ingestion pipeline (``docsearch/ingest.py``).
# ──────────────────────────────────────────────────────────────────────────
def _pack_settings() -> Tuple[int, int, float]:
    try:
        from django.conf import settings
        return (
            int(getattr(settings, "DOCSEARCH_CHUNK_TOKENS", 384)),
            int(getattr(settings, "DOCSEARCH_CHUNK_MAX_TOKENS", 512)),
            float(getattr(settings, "DOCSEARCH_CHUNK_OVERLAP", 0.15)),
        )
    except Exception:
        return 384, 512, 0.15


def _est_tokens(text: str) -> int:
    """Dependency-free token estimate: ~1.3 tokens per whitespace word."""
    return int(len((text or "").split()) * 1.3) + 1


def pack_passages(blocks: List[str], source: str, page) -> Tuple[List[str], List[Dict]]:
    """Pack the sentences of ``blocks`` (natural text blocks for ONE page/section,
    in reading order) into ~target-token passages with a small overlap.

    Returns ``(texts, metas)`` where each meta is ``{"source", "page"}`` — the
    only keys ``retrieval.build_chunks_dataframe`` consumes. Chunks are emitted
    contiguously and in order, which neighbour-expansion in ``DocIndex`` relies
    on.
    """
    target, hard_max, overlap = _pack_settings()
    sentences: List[str] = []
    for b in blocks:
        sentences.extend(_split_sentences(b))

    texts: List[str] = []
    metas: List[Dict] = []
    cur: List[str] = []
    cur_tok = 0

    def emit():
        nonlocal cur, cur_tok
        if not cur:
            return
        text = " ".join(cur).strip()
        if text:
            texts.append(text)
            metas.append({"source": source, "page": page})
        # Carry a small overlap tail into the next passage.
        keep = max(0, int(len(cur) * overlap))
        cur = cur[len(cur) - keep:] if keep else []
        cur_tok = _est_tokens(" ".join(cur)) if cur else 0

    for s in sentences:
        t = _est_tokens(s)
        if cur and (cur_tok + t > hard_max or cur_tok >= target):
            emit()
        cur.append(s)
        cur_tok += t

    final = " ".join(cur).strip()
    if final:
        texts.append(final)
        metas.append({"source": source, "page": page})
    return texts, metas


# ──────────────────────────────────────────────────────────────────────────
# DOCX  (the real LIMS knowledge corpus)
# ──────────────────────────────────────────────────────────────────────────
def _iter_docx_blocks(document):
    """Yield paragraphs and tables in document order (python-docx has no
    built-in body iterator that interleaves the two)."""
    from docx.document import Document as _Doc
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    body = document.element.body if isinstance(document, _Doc) else document._tc
    for child in body.iterchildren():
        tag = child.tag
        if tag.endswith("}p"):
            yield Paragraph(child, document)
        elif tag.endswith("}tbl"):
            yield Table(child, document)


def _extract_text_from_docx(path: Path) -> Tuple[List[str], List[Dict]]:
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(str(path))
    blocks_by_section: Dict[int, List[str]] = {}
    section = 1  # bumps at each heading, so "page" roughly maps to a section

    for block in _iter_docx_blocks(document):
        if isinstance(block, Paragraph):
            style = (block.style.name or "") if block.style else ""
            if style.startswith("Heading") or style == "Title":
                section += 1
            block_text = (block.text or "").strip()
        elif isinstance(block, Table):
            rows = []
            for row in block.rows:
                cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                if cells:
                    rows.append(" | ".join(dict.fromkeys(cells)))  # de-dup merged cells
            block_text = "\n".join(rows)
        else:
            continue

        if block_text:
            blocks_by_section.setdefault(section, []).append(block_text)

    texts: List[str] = []
    metas: List[Dict] = []
    for sec in sorted(blocks_by_section):
        t, m = pack_passages(blocks_by_section[sec], path.name, sec)
        texts.extend(t)
        metas.extend(m)
    return texts, metas


# ──────────────────────────────────────────────────────────────────────────
# PPTX
# ──────────────────────────────────────────────────────────────────────────
def _extract_text_from_pptx(path: Path) -> Tuple[List[str], List[Dict]]:
    from pptx import Presentation

    pres = Presentation(str(path))
    texts: List[str] = []
    metas: List[Dict] = []
    for slide_idx, slide in enumerate(pres.slides, start=1):
        parts = []
        for shape in slide.shapes:
            try:
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text:
                    parts.append(shape.text_frame.text.strip())
                elif getattr(shape, "text", None):
                    parts.append(str(shape.text).strip())
            except Exception:
                continue
        slide_text = "\n".join(p for p in parts if p).strip()
        if not slide_text:
            continue
        t, m = pack_passages([slide_text], path.name, slide_idx)
        texts.extend(t)
        metas.extend(m)
    return texts, metas


# ──────────────────────────────────────────────────────────────────────────
# PDF
# ──────────────────────────────────────────────────────────────────────────
def _extract_text_from_pdf(path: Path) -> Tuple[List[str], List[Dict]]:
    from pypdf import PdfReader

    texts: List[str] = []
    metas: List[Dict] = []
    reader = PdfReader(str(path))
    # Many corporate PDFs (calibration certs, ISO guidebooks, quality manuals) are
    # "encrypted" only to set permissions, with an EMPTY user password. Accessing
    # reader.pages then raises unless we decrypt first — which extract_file would
    # swallow, silently yielding 0 chunks. Decrypt with the empty password so the
    # document still indexes.
    if getattr(reader, "is_encrypted", False):
        try:
            reader.decrypt("")
        except Exception as exc:  # genuinely locked -> skip, don't abort the corpus
            print(f"⚠️  could not decrypt {path.name}: {exc}")
    for i, page in enumerate(reader.pages, start=1):
        try:
            page_text = page.extract_text() or ""
        except Exception as exc:  # one malformed page must not zero the whole doc
            print(f"⚠️  {path.name} page {i}: {exc}")
            continue
        if not page_text.strip():
            continue
        t, m = pack_passages([page_text], path.name, i)
        texts.extend(t)
        metas.extend(m)
    return texts, metas


# ──────────────────────────────────────────────────────────────────────────
# Plain text
# ──────────────────────────────────────────────────────────────────────────
def _extract_text_from_txt(path: Path) -> Tuple[List[str], List[Dict]]:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    blocks = [b.replace("\n", " ").strip() for b in re.split(r"\n\s*\n", raw)]  # paragraphs
    blocks = [b for b in blocks if b]
    return pack_passages(blocks, path.name, 1)


# ──────────────────────────────────────────────────────────────────────────
# Images (OCR -- optional, gated by settings.DOCSEARCH_ENABLE_OCR)
# ──────────────────────────────────────────────────────────────────────────
_ocr_reader = None


def _get_ocr_reader():
    global _ocr_reader
    if _ocr_reader is None:
        import easyocr  # heavy + downloads detector weights; imported lazily
        _ocr_reader = easyocr.Reader(["en"])
    return _ocr_reader


def _extract_text_from_image(path: Path) -> Tuple[List[str], List[Dict]]:
    try:
        from django.conf import settings
        if not getattr(settings, "DOCSEARCH_ENABLE_OCR", False):
            return [], []
    except Exception:
        return [], []
    try:
        lines = _get_ocr_reader().readtext(str(path), detail=0)
    except Exception as exc:  # easyocr missing or failed
        print(f"⚠️  OCR skipped for {path.name}: {exc}")
        return [], []
    ocr_text = "\n".join(l for l in lines if isinstance(l, str)).strip()
    if not ocr_text:
        return [], []
    return pack_passages([ocr_text], path.name, "image")


_EXTRACTORS = {
    ".docx": _extract_text_from_docx,
    ".pptx": _extract_text_from_pptx,
    ".pdf": _extract_text_from_pdf,
}


def extract_file(path: Path) -> Tuple[List[str], List[Dict]]:
    """Extract chunks from a single file. Returns ([], []) for unsupported or
    unreadable files (never raises)."""
    ext = path.suffix.lower()
    try:
        if ext in _EXTRACTORS:
            return _EXTRACTORS[ext](path)
        if ext in TEXT_EXTS:
            return _extract_text_from_txt(path)
        if ext in IMAGE_EXTS:
            return _extract_text_from_image(path)
    except Exception as exc:
        print(f"⚠️  Extraction failed for {path.name}: {exc}")
    return [], []


def extract_corpus(corpus_dir) -> Tuple[List[str], List[Dict]]:
    """Walk ``corpus_dir`` and extract chunks from every supported file."""
    corpus = Path(corpus_dir)
    corpus.mkdir(parents=True, exist_ok=True)
    chunk_texts: List[str] = []
    chunk_metas: List[Dict] = []
    for path in sorted(corpus.iterdir()):
        if not path.is_file() or path.suffix.lower() not in SUPPORTED_EXTS:
            continue
        texts, metas = extract_file(path)
        chunk_texts.extend(texts)
        chunk_metas.extend(metas)
    print(f"✓ extract_corpus: {len(chunk_texts)} chunks from {corpus}")
    return chunk_texts, chunk_metas
